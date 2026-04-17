#!/usr/bin/env python3
"""
PPT generator script using python-pptx.
Reads slide_plan.json from the current directory and produces output.pptx.

slide_plan.json schema:
{
  "title": "Presentation title",
  "author": "Author name (optional)",
  "theme": "business|dark|minimal",
  "slides": [
    {
      "type": "title|section|content|closing",
      "title": "Slide title",
      "key_message": "One-line core message",
      "bullets": ["bullet 1", "bullet 2"],
      "notes": "Presenter notes",
      "visual_hint": "Description of recommended visual"
    }
  ]
}
"""
import json
import sys
import subprocess

def ensure_pptx():
    try:
        import pptx
    except ImportError:
        print("Installing python-pptx...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])

ensure_pptx()

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Theme palettes ─────────────────────────────────────────────────────────────
THEMES = {
    "business": {
        "bg":        RGBColor(0xFF, 0xFF, 0xFF),
        "accent":    RGBColor(0x1F, 0x4E, 0x79),   # dark navy
        "accent2":   RGBColor(0x2E, 0x75, 0xB6),   # medium blue
        "title_fg":  RGBColor(0xFF, 0xFF, 0xFF),
        "body_fg":   RGBColor(0x1F, 0x1F, 0x1F),
        "bullet_fg": RGBColor(0x2E, 0x75, 0xB6),
        "note_fg":   RGBColor(0x59, 0x59, 0x59),
    },
    "dark": {
        "bg":        RGBColor(0x1A, 0x1A, 0x2E),
        "accent":    RGBColor(0x16, 0x21, 0x3E),
        "accent2":   RGBColor(0x0F, 0x3D, 0x60),
        "title_fg":  RGBColor(0xE9, 0xF1, 0xFA),
        "body_fg":   RGBColor(0xC8, 0xD8, 0xE8),
        "bullet_fg": RGBColor(0x00, 0xB4, 0xD8),
        "note_fg":   RGBColor(0x90, 0xA0, 0xB0),
    },
    "minimal": {
        "bg":        RGBColor(0xFA, 0xFA, 0xFA),
        "accent":    RGBColor(0x2D, 0x2D, 0x2D),
        "accent2":   RGBColor(0x55, 0x55, 0x55),
        "title_fg":  RGBColor(0xFF, 0xFF, 0xFF),
        "body_fg":   RGBColor(0x2D, 0x2D, 0x2D),
        "bullet_fg": RGBColor(0x55, 0x55, 0x55),
        "note_fg":   RGBColor(0x77, 0x77, 0x77),
    },
}

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=RGBColor(0,0,0),
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def make_title_slide(prs, slide_data, theme):
    t = THEMES[theme]
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_bg(slide, t["bg"])

    # Full accent bar
    add_rect(slide, 0, 0, W, H * 0.65, t["accent"])
    # Bottom accent strip
    add_rect(slide, 0, H * 0.65, W, Inches(0.12), t["accent2"])

    title = slide_data.get("title", "")
    key_msg = slide_data.get("key_message", "")
    notes_text = slide_data.get("notes", "")

    add_text_box(slide, title,
                 Inches(1), Inches(1.5), Inches(11), Inches(2.2),
                 font_size=40, bold=True, color=t["title_fg"], align=PP_ALIGN.CENTER)

    if key_msg:
        add_text_box(slide, key_msg,
                     Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
                     font_size=20, bold=False, color=RGBColor(0xCC, 0xDD, 0xEE),
                     align=PP_ALIGN.CENTER)

    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    return slide


def make_section_slide(prs, slide_data, theme):
    t = THEMES[theme]
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, t["accent"])

    title = slide_data.get("title", "")
    key_msg = slide_data.get("key_message", "")

    add_text_box(slide, title,
                 Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
                 font_size=36, bold=True, color=t["title_fg"], align=PP_ALIGN.CENTER)

    if key_msg:
        add_text_box(slide, key_msg,
                     Inches(1.5), Inches(4.2), Inches(10), Inches(0.8),
                     font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE),
                     align=PP_ALIGN.CENTER)

    notes_text = slide_data.get("notes", "")
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    return slide


def make_content_slide(prs, slide_data, theme):
    t = THEMES[theme]
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, t["bg"])

    # Top title bar
    add_rect(slide, 0, 0, W, Inches(1.2), t["accent"])

    title = slide_data.get("title", "")
    key_msg = slide_data.get("key_message", "")
    bullets = slide_data.get("bullets", [])
    visual_hint = slide_data.get("visual_hint", "")
    notes_text = slide_data.get("notes", "")

    add_text_box(slide, title,
                 Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.9),
                 font_size=24, bold=True, color=t["title_fg"])

    if key_msg:
        add_text_box(slide, f"✦  {key_msg}",
                     Inches(0.3), Inches(1.3), Inches(12.7), Inches(0.6),
                     font_size=14, bold=True, color=t["accent2"])

    # Bullets
    if bullets:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = f"•  {bullet}"
            run.font.size = Pt(16)
            run.font.color.rgb = t["body_fg"]

    # Visual hint as a small label at bottom-right
    if visual_hint:
        add_text_box(slide, f"[시각 자료: {visual_hint}]",
                     Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6),
                     font_size=10, color=t["note_fg"], align=PP_ALIGN.RIGHT)

    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    return slide


def make_closing_slide(prs, slide_data, theme):
    t = THEMES[theme]
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, t["accent"])

    title = slide_data.get("title", "감사합니다")
    key_msg = slide_data.get("key_message", "")
    bullets = slide_data.get("bullets", [])

    add_text_box(slide, title,
                 Inches(1), Inches(2.2), Inches(11.3), Inches(1.5),
                 font_size=40, bold=True, color=t["title_fg"], align=PP_ALIGN.CENTER)

    content_parts = []
    if key_msg:
        content_parts.append(key_msg)
    content_parts.extend(bullets)

    if content_parts:
        combined = "\n".join(content_parts)
        add_text_box(slide, combined,
                     Inches(1.5), Inches(3.9), Inches(10), Inches(2.5),
                     font_size=16, color=RGBColor(0xCC, 0xDD, 0xEE),
                     align=PP_ALIGN.CENTER)

    notes_text = slide_data.get("notes", "")
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    return slide


SLIDE_MAKERS = {
    "title":   make_title_slide,
    "section": make_section_slide,
    "content": make_content_slide,
    "closing": make_closing_slide,
}


def build_presentation(plan: dict, output_path: str):
    theme = plan.get("theme", "business")
    if theme not in THEMES:
        theme = "business"

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    for slide_data in plan.get("slides", []):
        slide_type = slide_data.get("type", "content")
        maker = SLIDE_MAKERS.get(slide_type, make_content_slide)
        maker(prs, slide_data, theme)

    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    plan_path = sys.argv[1] if len(sys.argv) > 1 else "slide_plan.json"
    out_path   = sys.argv[2] if len(sys.argv) > 2 else "presentation.pptx"

    with open(plan_path, encoding="utf-8") as f:
        plan = json.load(f)

    build_presentation(plan, out_path)
